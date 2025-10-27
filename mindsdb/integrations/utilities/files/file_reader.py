from dataclasses import dataclass, astuple
import traceback
import json
import csv
from io import BytesIO, StringIO, IOBase
from pathlib import Path
import codecs
from typing import List, Generator

import filetype
import pandas as pd
from charset_normalizer import from_bytes
from mindsdb.interfaces.knowledge_base.preprocessing.text_splitter import TextSplitter

from mindsdb.utilities import log

logger = log.getLogger(__name__)

DEFAULT_CHUNK_SIZE = 800
DEFAULT_CHUNK_OVERLAP = 80


class FileProcessingError(Exception): ...


@dataclass(frozen=True, slots=True)
class _SINGLE_PAGE_FORMAT:
    CSV: str = "csv"
    JSON: str = "json"
    TXT: str = "txt"
    PDF: str = "pdf"
    MD: str = "md"
    DOC: str = "doc"
    DOCX: str = "docx"
    PPTX: str = "pptx"
    PPT: str = "ppt"
    PARQUET: str = "parquet"


SINGLE_PAGE_FORMAT = _SINGLE_PAGE_FORMAT()


@dataclass(frozen=True, slots=True)
class _MULTI_PAGE_FORMAT:
    XLSX: str = "xlsx"


MULTI_PAGE_FORMAT = _MULTI_PAGE_FORMAT()


def decode(file_obj: IOBase) -> StringIO:
    file_obj.seek(0)
    byte_str = file_obj.read()
    # Move it to StringIO
    try:
        # Handle Microsoft's BOM "special" UTF-8 encoding
        if byte_str.startswith(codecs.BOM_UTF8):
            data_str = StringIO(byte_str.decode("utf-8-sig"))
        else:
            file_encoding_meta = from_bytes(
                byte_str[: 32 * 1024],
                steps=32,  # Number of steps/block to extract from my_byte_str
                chunk_size=1024,  # Set block size of each extraction)
                explain=False,
            )
            best_meta = file_encoding_meta.best()
            errors = "strict"
            if best_meta is not None:
                encoding = file_encoding_meta.best().encoding

                try:
                    data_str = StringIO(byte_str.decode(encoding, errors))
                except UnicodeDecodeError:
                    encoding = "utf-8"
                    errors = "replace"

                    data_str = StringIO(byte_str.decode(encoding, errors))
            else:
                encoding = "utf-8"
                errors = "replace"

                data_str = StringIO(byte_str.decode(encoding, errors))
    except Exception as e:
        logger.error(traceback.format_exc())
        raise FileProcessingError("Could not load into string") from e

    return data_str


class FormatDetector:
    supported_formats = astuple(SINGLE_PAGE_FORMAT) + astuple(MULTI_PAGE_FORMAT)
    multipage_formats = astuple(MULTI_PAGE_FORMAT)

    def __init__(
        self,
        path: str | None = None,
        name: str | None = None,
        file: IOBase | None = None,
    ):
        """
        File format detector
        One of these arguments has to be passed: `path` or `file`

        :param path: path to the file
        :param name: name of the file
        :param file: file descriptor (via open(...), of BytesIO(...))
        """
        if path is not None:
            file = open(path, "rb")

        elif file is not None:
            if name is None:
                if hasattr(file, "name"):
                    path = file.name
                else:
                    path = "file"
        else:
            raise FileProcessingError("Wrong arguments: path or file is required")

        if name is None:
            name = Path(path).name

        self.name = name
        self.file_obj = file
        self.format = None

        self.parameters = {}

    def get_format(self) -> str:
        if self.format is not None:
            return self.format

        format = self.get_format_by_name()
        if format is not None:
            if format not in self.supported_formats:
                raise FileProcessingError(f"Not supported format: {format}")

        if format is None and self.file_obj is not None:
            format = self.get_format_by_content()
            self.file_obj.seek(0)

        if format is None:
            raise FileProcessingError(f"Unable to detect format: {self.name}")

        self.format = format
        return format

    def get_format_by_name(self):
        extension = Path(self.name).suffix.strip(".").lower()
        if extension == "tsv":
            extension = "csv"
            self.parameters["delimiter"] = "\t"
        elif extension == "markdown":
            extension = "md"

        return extension or None

    def get_format_by_content(self):
        if self.is_parquet(self.file_obj):
            return SINGLE_PAGE_FORMAT.PARQUET

        file_type = filetype.guess(self.file_obj)
        if file_type is not None:
            if file_type.mime in {
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            }:
                return MULTI_PAGE_FORMAT.XLSX

            if file_type.mime == "application/pdf":
                return SINGLE_PAGE_FORMAT.PDF

            if file_type.mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return SINGLE_PAGE_FORMAT.DOCX

            if file_type.mime == "application/msword":
                return SINGLE_PAGE_FORMAT.DOC

            if file_type.mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return SINGLE_PAGE_FORMAT.PPTX

            if file_type.mime == "application/vnd.ms-powerpoint":
                return SINGLE_PAGE_FORMAT.PPT

        file_obj = decode(self.file_obj)

        if self.is_json(file_obj):
            return SINGLE_PAGE_FORMAT.JSON

        if self.is_csv(file_obj):
            return SINGLE_PAGE_FORMAT.CSV

    @staticmethod
    def is_json(data_obj: StringIO) -> bool:
        # see if its JSON
        data_obj.seek(0)
        text = data_obj.read(100).strip()
        data_obj.seek(0)
        if len(text) > 0:
            # it looks like a json, then try to parse it
            if text.startswith("{") or text.startswith("["):
                try:
                    json.loads(data_obj.read())
                    return True
                except Exception:
                    return False
        return False

    @classmethod
    def is_csv(cls, data_obj: StringIO) -> bool:
        data_obj.seek(0)
        sample = data_obj.readline()  # trying to get dialect from header
        try:
            data_obj.seek(0)
            csv.Sniffer().sniff(sample)

            # Avoid a false-positive for json files
            if cls.is_json(data_obj):
                return False
            return True
        except Exception:
            return False

    @staticmethod
    def is_parquet(data: IOBase) -> bool:
        # Check first and last 4 bytes equal to PAR1.
        # Refer: https://parquet.apache.org/docs/file-format/
        parquet_sig = b"PAR1"
        data.seek(0, 0)
        start_meta = data.read(4)
        data.seek(-4, 2)
        end_meta = data.read()
        data.seek(0)
        if start_meta == parquet_sig and end_meta == parquet_sig:
            return True
        return False


def format_column_names(df: pd.DataFrame):
    df.columns = [column.strip(" \t") for column in df.columns]
    if len(df.columns) != len(set(df.columns)) or any(len(column_name) == 0 for column_name in df.columns):
        raise FileProcessingError("Each column should have a unique and non-empty name.")


class FileReader(FormatDetector):
    def _get_fnc(self):
        format = self.get_format()
        func = getattr(self, f"read_{format}", None)
        if func is None:
            raise FileProcessingError(f"Unsupported format: {format}")

        if format in astuple(MULTI_PAGE_FORMAT):

            def format_multipage(*args, **kwargs):
                for page_number, df in func(*args, **kwargs):
                    format_column_names(df)
                    yield page_number, df

            return format_multipage

        def format_singlepage(*args, **kwargs) -> pd.DataFrame:
            """Check that the columns have unique not-empty names"""
            df = func(*args, **kwargs)
            format_column_names(df)
            return df

        return format_singlepage

    def get_pages(self, **kwargs) -> List[str]:
        """
        Get list of tables in file
        """
        format = self.get_format()
        if format not in self.multipage_formats:
            # only one table
            return ["main"]

        func = self._get_fnc()
        self.file_obj.seek(0)

        return [name for name, _ in func(self.file_obj, only_names=True, **kwargs)]

    def get_contents(self, **kwargs) -> dict[str, pd.DataFrame]:
        """
        Get all info(pages with content) from file as dict: {tablename, content}
        """
        func = self._get_fnc()
        self.file_obj.seek(0)

        format = self.get_format()
        if format not in self.multipage_formats:
            # only one table
            return {"main": func(self.file_obj, name=self.name, **kwargs)}

        return {name: df for name, df in func(self.file_obj, **kwargs)}

    def get_page_content(self, page_name: str | None = None, **kwargs) -> pd.DataFrame:
        """
        Get content of a single table
        """
        func = self._get_fnc()
        self.file_obj.seek(0)

        format = self.get_format()
        if format not in self.multipage_formats:
            # only one table
            return func(self.file_obj, name=self.name, **kwargs)

        for _, df in func(self.file_obj, name=self.name, page_name=page_name, **kwargs):
            return df

    @staticmethod
    def _get_csv_dialect(buffer, delimiter: str | None = None) -> csv.Dialect | None:
        sample = buffer.readline()  # trying to get dialect from header
        buffer.seek(0)
        try:
            if isinstance(sample, bytes):
                sample = sample.decode()

            if delimiter is not None:
                accepted_csv_delimiters = [delimiter]
            else:
                accepted_csv_delimiters = [",", "\t", ";"]
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=accepted_csv_delimiters)
                dialect.doublequote = True  # assume that all csvs have " as string escape
            except Exception:
                dialect = csv.reader(sample).dialect
                if dialect.delimiter not in accepted_csv_delimiters:
                    raise FileProcessingError(f"CSV delimeter '{dialect.delimiter}' is not supported")

        except csv.Error:
            dialect = None
        return dialect

    @classmethod
    def read_csv(cls, file_obj: BytesIO, delimiter: str | None = None, **kwargs) -> pd.DataFrame:
        file_obj = decode(file_obj)
        dialect = cls._get_csv_dialect(file_obj, delimiter=delimiter)
        return pd.read_csv(file_obj, sep=dialect.delimiter, index_col=False)

    @staticmethod
    def read_txt(file_obj: BytesIO, name: str | None = None, chunk_size: int | None = None, chunk_overlap: int | None = None, **kwargs) -> pd.DataFrame:
        # the lib is heavy, so import it only when needed

        file_obj = decode(file_obj)

        text = file_obj.read()

        # If chunk_size is 0 or negative, return full text without chunking
        if chunk_size is not None and chunk_size <= 0:
            return pd.DataFrame([{"content": text, "metadata": {"source_file": name, "file_format": "txt"}}])

        # Use provided chunk_size and chunk_overlap if available, otherwise use defaults
        _chunk_size = chunk_size if chunk_size is not None else DEFAULT_CHUNK_SIZE
        _chunk_overlap = chunk_overlap if chunk_overlap is not None else DEFAULT_CHUNK_OVERLAP

        text_splitter = TextSplitter(chunk_size=_chunk_size, chunk_overlap=_chunk_overlap)

        docs = text_splitter.split_text(text)
        return pd.DataFrame([{"content": doc, "metadata": {"source_file": name, "file_format": "txt"}} for doc in docs])

    @staticmethod
    def read_md(file_obj: BytesIO, name: str | None = None, chunk_size: int | None = None, chunk_overlap: int | None = None, **kwargs) -> pd.DataFrame:
        # Read markdown files similar to txt

        file_obj = decode(file_obj)

        text = file_obj.read()

        # If chunk_size is 0 or negative, return full text without chunking
        if chunk_size is not None and chunk_size <= 0:
            return pd.DataFrame([{"content": text, "metadata": {"source_file": name, "file_format": "md"}}])

        # Use provided chunk_size and chunk_overlap if available, otherwise use defaults
        _chunk_size = chunk_size if chunk_size is not None else DEFAULT_CHUNK_SIZE
        _chunk_overlap = chunk_overlap if chunk_overlap is not None else DEFAULT_CHUNK_OVERLAP

        text_splitter = TextSplitter(chunk_size=_chunk_size, chunk_overlap=_chunk_overlap)

        docs = text_splitter.split_text(text)
        return pd.DataFrame([{"content": doc, "metadata": {"source_file": name, "file_format": "md"}} for doc in docs])

    @staticmethod
    def read_docx(file_obj: BytesIO, name: str | None = None, chunk_size: int | None = None, chunk_overlap: int | None = None, **kwargs) -> pd.DataFrame:
        # the lib is heavy, so import it only when needed
        try:
            import docx
            from docx import Document as DocxDocument
        except ImportError as e:
            raise FileProcessingError(
                "python-docx package is required to read DOCX files. "
                "Install it with: pip install python-docx"
            ) from e

        doc = DocxDocument(file_obj)

        # Extract text from all paragraphs
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])

        # If chunk_size is 0 or negative, return full text without chunking
        if chunk_size is not None and chunk_size <= 0:
            return pd.DataFrame([{"content": text, "metadata": {"source_file": name, "file_format": "docx"}}])

        # Use provided chunk_size and chunk_overlap if available, otherwise use defaults
        _chunk_size = chunk_size if chunk_size is not None else DEFAULT_CHUNK_SIZE
        _chunk_overlap = chunk_overlap if chunk_overlap is not None else DEFAULT_CHUNK_OVERLAP

        text_splitter = TextSplitter(chunk_size=_chunk_size, chunk_overlap=_chunk_overlap)

        docs = text_splitter.split_text(text)
        return pd.DataFrame([{"content": doc, "metadata": {"source_file": name, "file_format": "docx"}} for doc in docs])

    @staticmethod
    def read_doc(file_obj: BytesIO, name: str | None = None, chunk_size: int | None = None, chunk_overlap: int | None = None, **kwargs) -> pd.DataFrame:
        # DOC files (older Word format) are more complex and require different handling
        # For now, we'll try to use python-docx which sometimes works with .doc files
        # or fall back to treating as text
        try:
            import docx
            from docx import Document as DocxDocument
            doc = DocxDocument(file_obj)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
        except ImportError:
            raise FileProcessingError(
                "python-docx package is required to read DOC files. "
                "Install it with: pip install python-docx"
            )
        except Exception:
            # If docx fails, try reading as text
            file_obj = decode(file_obj)
            text = file_obj.read()

        # If chunk_size is 0 or negative, return full text without chunking
        if chunk_size is not None and chunk_size <= 0:
            return pd.DataFrame([{"content": text, "metadata": {"source_file": name, "file_format": "doc"}}])

        # Use provided chunk_size and chunk_overlap if available, otherwise use defaults
        _chunk_size = chunk_size if chunk_size is not None else DEFAULT_CHUNK_SIZE
        _chunk_overlap = chunk_overlap if chunk_overlap is not None else DEFAULT_CHUNK_OVERLAP

        text_splitter = TextSplitter(chunk_size=_chunk_size, chunk_overlap=_chunk_overlap)

        docs = text_splitter.split_text(text)
        return pd.DataFrame([{"content": doc, "metadata": {"source_file": name, "file_format": "doc"}} for doc in docs])

    @staticmethod
    def read_pptx(file_obj: BytesIO, name: str | None = None, chunk_size: int | None = None, chunk_overlap: int | None = None, **kwargs) -> pd.DataFrame:
        # the lib is heavy, so import it only when needed
        try:
            from pptx import Presentation
        except ImportError as e:
            raise FileProcessingError(
                "python-pptx package is required to read PPTX files. "
                "Install it with: pip install python-pptx"
            ) from e

        prs = Presentation(file_obj)

        # Extract text from all slides
        text_parts = []
        for slide in prs.slides:
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

                # Extract text from tables if present
                if shape.shape_type == 19:  # Table shape type
                    try:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text_parts.append(cell.text)
                    except Exception:
                        pass

            # Extract notes if present
            if slide.has_notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        text_parts.append(notes_text)
                except Exception:
                    pass

        text = "\n".join(text_parts)

        # If chunk_size is 0 or negative, return full text without chunking
        if chunk_size is not None and chunk_size <= 0:
            return pd.DataFrame([{"content": text, "metadata": {"source_file": name, "file_format": "pptx"}}])

        # Use provided chunk_size and chunk_overlap if available, otherwise use defaults
        _chunk_size = chunk_size if chunk_size is not None else DEFAULT_CHUNK_SIZE
        _chunk_overlap = chunk_overlap if chunk_overlap is not None else DEFAULT_CHUNK_OVERLAP

        text_splitter = TextSplitter(chunk_size=_chunk_size, chunk_overlap=_chunk_overlap)

        docs = text_splitter.split_text(text)
        return pd.DataFrame([{"content": doc, "metadata": {"source_file": name, "file_format": "pptx"}} for doc in docs])

    @staticmethod
    def read_ppt(file_obj: BytesIO, name: str | None = None, chunk_size: int | None = None, chunk_overlap: int | None = None, **kwargs) -> pd.DataFrame:
        # PPT files (older PowerPoint format) are binary and more complex
        # Try to use python-pptx which sometimes works with .ppt files
        try:
            from pptx import Presentation
        except ImportError as e:
            raise FileProcessingError(
                "python-pptx package is required to read PPT files. "
                "Install it with: pip install python-pptx"
            ) from e

        try:
            prs = Presentation(file_obj)

            # Extract text from all slides (same logic as PPTX)
            text_parts = []
            for slide in prs.slides:
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                    # Extract text from tables if present
                    if shape.shape_type == 19:  # Table shape type
                        try:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        text_parts.append(cell.text)
                        except Exception:
                            pass

                # Extract notes if present
                if slide.has_notes_slide:
                    try:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text.strip():
                            text_parts.append(notes_text)
                    except Exception:
                        pass

            text = "\n".join(text_parts)
        except Exception as e:
            # If python-pptx fails with PPT format, raise an informative error
            raise FileProcessingError(
                f"Unable to read PPT file '{name}'. The PPT format (older PowerPoint) "
                "is a complex binary format. Consider converting it to PPTX format first. "
                f"Error: {str(e)}"
            ) from e

        # If chunk_size is 0 or negative, return full text without chunking
        if chunk_size is not None and chunk_size <= 0:
            return pd.DataFrame([{"content": text, "metadata": {"source_file": name, "file_format": "ppt"}}])

        # Use provided chunk_size and chunk_overlap if available, otherwise use defaults
        _chunk_size = chunk_size if chunk_size is not None else DEFAULT_CHUNK_SIZE
        _chunk_overlap = chunk_overlap if chunk_overlap is not None else DEFAULT_CHUNK_OVERLAP

        text_splitter = TextSplitter(chunk_size=_chunk_size, chunk_overlap=_chunk_overlap)

        docs = text_splitter.split_text(text)
        return pd.DataFrame([{"content": doc, "metadata": {"source_file": name, "file_format": "ppt"}} for doc in docs])

    @staticmethod
    def read_pdf(file_obj: BytesIO, name: str | None = None, chunk_size: int | None = None, chunk_overlap: int | None = None, **kwargs) -> pd.DataFrame:
        # the libs are heavy, so import it only when needed
        import fitz  # pymupdf

        with fitz.open(stream=file_obj.read()) as pdf:  # open pdf
            text = '\n'.join([page.get_text() for page in pdf])

        # If chunk_size is 0 or negative, return full text without chunking
        if chunk_size is not None and chunk_size <= 0:
            return pd.DataFrame(
                {
                    "content": [text],
                    "metadata": [{"file_format": "pdf", "source_file": name, "chunk_size": 0, "chunk_overlap": 0}],
                }
            )

        # Use provided chunk_size and chunk_overlap if available, otherwise use defaults
        _chunk_size = chunk_size if chunk_size is not None else DEFAULT_CHUNK_SIZE
        _chunk_overlap = chunk_overlap if chunk_overlap is not None else DEFAULT_CHUNK_OVERLAP

        text_splitter = TextSplitter(chunk_size=_chunk_size, chunk_overlap=_chunk_overlap)

        split_text = text_splitter.split_text(text)

        return pd.DataFrame(
            {
                "content": split_text,
                "metadata": [{"file_format": "pdf", "source_file": name, "chunk_size": _chunk_size, "chunk_overlap": _chunk_overlap}] * len(split_text),
            }
        )

    @staticmethod
    def read_json(file_obj: BytesIO, **kwargs) -> pd.DataFrame:
        file_obj = decode(file_obj)
        file_obj.seek(0)
        json_doc = json.loads(file_obj.read())
        return pd.json_normalize(json_doc, max_level=0)

    @staticmethod
    def read_parquet(file_obj: BytesIO, **kwargs) -> pd.DataFrame:
        return pd.read_parquet(file_obj)

    @staticmethod
    def read_xlsx(
        file_obj: BytesIO,
        page_name: str | None = None,
        only_names: bool = False,
        **kwargs,
    ) -> Generator[tuple[str, pd.DataFrame | None], None, None]:
        with pd.ExcelFile(file_obj) as xls:
            if page_name is not None:
                # return specific page
                yield page_name, pd.read_excel(xls, sheet_name=page_name)

            for page_name in xls.sheet_names:
                if only_names:
                    # extract only pages names
                    df = None
                else:
                    df = pd.read_excel(xls, sheet_name=page_name)
                yield page_name, df
